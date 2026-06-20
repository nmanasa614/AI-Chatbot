from pptx import Presentation
import os

def create_ppt(topic):

    prs = Presentation()

    slides_content = [
        ("Title", topic),
        ("Introduction", f"Introduction to {topic}"),
        ("Key Concepts", f"Important concepts of {topic}"),
        ("Applications", f"Applications of {topic}"),
        ("Advantages", f"Advantages of {topic}"),
        ("Challenges", f"Challenges of {topic}"),
        ("Conclusion", f"Summary of {topic}")
    ]

    for title, content in slides_content:

        slide = prs.slides.add_slide(
            prs.slide_layouts[1]
        )

        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    downloads = os.path.join(
        os.path.expanduser("~"),
        "Downloads"
    )

    filename = os.path.join(
        downloads,
        "generated_presentation.pptx"
    )

    prs.save(filename)

    return filename